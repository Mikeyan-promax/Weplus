"""
文档服务
处理文档上传、解析、分块和向量化
"""

import os
import io
import logging
from typing import List, Dict, Any, Optional, BinaryIO
from datetime import datetime
import asyncio
import json

# 文档处理库
import PyPDF2
from docx import Document as DocxDocument

# 本地服务导入
from .rag_service import rag_service
from .postgresql_vector_service import postgresql_vector_service
from ..models.document import Document, DocumentChunk, calculate_file_hash, get_file_type, validate_file_type

# 配置日志
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class DocumentService:
    """文档服务类"""
    
    def __init__(self):
        # 配置参数
        self.upload_dir = "data/uploads"
        self.max_file_size = 50 * 1024 * 1024  # 50MB
        self.supported_types = ['pdf', 'word', 'text', 'markdown', 'html', 'powerpoint', 'excel']
        
        # 确保上传目录存在
        os.makedirs(self.upload_dir, exist_ok=True)
        
        logger.info("文档服务初始化完成")

    # ==================== 文档基础操作 ====================
    
    async def upload_document(
        self, 
        file_content: bytes, 
        filename: str, 
        metadata: Dict[str, Any] = None
    ) -> Dict[str, Any]:
        """
        上传并处理文档
        
        函数说明（中文）：
        - 执行严格的文件校验（大小、类型、扩展名匹配策略）。
        - 使用内容哈希检测重复；若存在失败/半成品记录，先清理再继续。
        - 先解析文本，确保可向量化后再创建数据库记录，避免“上传失败但仍保存”。
        - 若后续流程任一步失败，执行彻底清理（删除数据库记录、向量数据、本地文件）。
        
        返回结构：
        - 成功：{"success": True, "document_id": int, ...}
        - 失败：{"success": False, "error": str, ...}
        
        Args:
            file_content: 文件内容
            filename: 文件名
            metadata: 文档元数据
            
        Returns:
            处理结果
        """
        try:
            # 验证文件
            validation_result = self._validate_file(file_content, filename)
            if not validation_result["valid"]:
                return {
                    "success": False,
                    "error": validation_result["error"],
                    "timestamp": datetime.now().isoformat()
                }
            
            # 计算文件哈希
            content_hash = calculate_file_hash(file_content)
            file_type = get_file_type(filename)
            file_size = len(file_content)
            
            # 检查文档是否已存在；若为失败/半成品状态则清理后继续
            existing_doc = Document.get_by_hash(content_hash)
            if existing_doc:
                if existing_doc.status in ("failed", "uploaded", "processing"):
                    try:
                        # 清理旧的失败或半成品记录
                        await postgresql_vector_service.delete_document(str(existing_doc.id))
                    except Exception:
                        pass
                    try:
                        # 删除本地同哈希文件
                        import glob
                        for f in glob.glob(os.path.join(self.upload_dir, f"{existing_doc.content_hash}.*")):
                            try:
                                os.remove(f)
                            except Exception:
                                pass
                    except Exception:
                        pass
                    try:
                        existing_doc.delete()
                    except Exception:
                        pass
                else:
                    return {
                        "success": False,
                        "error": "文档已存在",
                        "document_id": existing_doc.id,
                        "timestamp": datetime.now().isoformat()
                    }
            
            # 先解析文本，确认可向量化；避免创建失败记录
            text_content = await self._extract_text(file_content, file_type)
            if not text_content:
                return {
                    "success": False,
                    "error": "文档内容解析失败",
                    "timestamp": datetime.now().isoformat()
                }
            
            # 文档分块（解析后，未写数据库，避免失败残留）
            chunks = rag_service.chunk_text(text_content)
            
            # 创建文档记录
            document = None
            try:
                document = Document.create_document(
                    filename=filename,
                    file_type=file_type,
                    file_size=file_size,
                    content_hash=content_hash,
                    metadata=metadata or {}
                )
            except Exception as e:
                logger.error(f"创建文档记录失败: {str(e)}")
                return {
                    "success": False,
                    "error": "创建文档记录失败",
                    "timestamp": datetime.now().isoformat()
                }
            
            # 保存文件到本地并进入处理状态
            file_path = await self._save_file(file_content, filename, content_hash)
            document.update_status("processing")
            
            # 保存分块到数据库（已改为使用PostgreSQL向量存储，不再重复写入document_chunks记录）
            # 说明：避免与向量服务写入同一表导致唯一键冲突，这里跳过数据库分块写入，仅依赖向量存储中的document_chunks。
            
            # 向量化并存储
            try:
                vector_result = await postgresql_vector_service.add_document(
                    document_id=str(document.id),
                    chunks=chunks,
                    metadata={
                        "filename": filename,
                        "file_type": file_type,
                        "file_size": file_size,
                        "upload_time": datetime.now().isoformat(),
                        **(metadata or {})
                    }
                )
            except Exception as e:
                await self._cleanup_failed_upload(document)
                return {
                    "success": False,
                    "error": f"向量化或存储失败: {str(e)}",
                    "timestamp": datetime.now().isoformat()
                }
            
            # 更新文档状态
            document.update_status("processed")
            
            result = {
                "success": True,
                "document_id": document.id,
                "filename": filename,
                "file_type": file_type,
                "file_size": file_size,
                "chunks_count": len(chunks),
                "content_length": len(text_content),
                "vector_result": vector_result,
                "timestamp": datetime.now().isoformat()
            }
            
            logger.info(f"文档 {filename} 上传处理成功")
            return result
            
        except Exception as e:
            logger.error(f"文档上传处理失败: {str(e)}")
            return {
                "success": False,
                "error": f"处理失败: {str(e)}",
                "timestamp": datetime.now().isoformat()
            }
    
    async def get_document_list(
        self, 
        offset: int = 0, 
        limit: int = 20,
        category_id: Optional[int] = None,
        search: Optional[str] = None,
        status: Optional[str] = None,
        sort_by: str = "upload_time",
        sort_order: str = "desc"
    ) -> Dict[str, Any]:
        """获取文档列表（支持筛选和搜索）"""
        try:
            # 使用Document.get_all方法获取文档列表
            documents = Document.get_all(
                limit=limit, 
                offset=offset, 
                category_id=category_id,
                status=status
            )
            
            # 如果有搜索条件，使用搜索方法
            if search:
                documents = Document.search(search, limit=limit, offset=offset)
            
            document_list = []
            for doc in documents:
                doc_dict = doc.to_dict()
                
                # 添加分类信息
                if doc.category_id:
                    from app.models.document import Category
                    category = Category.get_by_id(doc.category_id)
                    doc_dict["category"] = category.to_dict() if category else None
                
                document_list.append(doc_dict)
            
            return {
                "documents": document_list,
                "total": len(document_list),
                "offset": offset,
                "limit": limit
            }
        except Exception as e:
            logger.error(f"获取文档列表失败: {str(e)}")
            # 移除HTTPException依赖，直接返回错误信息
            return {
                "documents": [],
                "total": 0,
                "offset": offset,
                "limit": limit,
                "error": str(e)
            }
    
    async def get_document_detail(self, document_id: int) -> Dict[str, Any]:
        """获取文档详情"""
        try:
            document = Document.get_by_id(document_id)
            if not document:
                return {
                    "success": False,
                    "error": "文档不存在"
                }
            
            # 获取文档分块
            chunks = DocumentChunk.get_by_document_id(document_id)
            
            result = document.to_dict()
            result.update({
                "chunks": [chunk.to_dict() for chunk in chunks],
                "chunks_count": len(chunks),
                "success": True
            })
            
            # 添加分类和标签详细信息
            if hasattr(document, 'category_id') and document.category_id:
                result["category"] = self._get_category_info(document.category_id)
            if hasattr(document, 'tag_ids') and document.tag_ids:
                result["tags"] = self._get_tags_info(document.tag_ids)
            
            return result
            
        except Exception as e:
            logger.error(f"获取文档详情失败: {str(e)}")
            return {
                "success": False,
                "error": str(e)
            }
    
    async def update_document(
        self,
        document_id: int,
        title: Optional[str] = None,
        description: Optional[str] = None,
        category_id: Optional[int] = None,
        tag_ids: Optional[List[int]] = None
    ) -> Dict[str, Any]:
        """更新文档信息"""
        try:
            document = Document.get_by_id(document_id)
            if not document:
                return {
                    "success": False,
                    "error": "文档不存在"
                }
            
            # 更新文档信息
            update_data = {}
            if title is not None:
                update_data["title"] = title
            if description is not None:
                update_data["description"] = description
            if category_id is not None:
                update_data["category_id"] = category_id
            if tag_ids is not None:
                update_data["tag_ids"] = tag_ids
            
            success = document.update_info(update_data)
            
            if success:
                return {
                    "success": True,
                    "document_id": document_id,
                    "updated_fields": list(update_data.keys()),
                    "timestamp": datetime.now().isoformat()
                }
            else:
                return {
                    "success": False,
                    "error": "更新失败"
                }
            
        except Exception as e:
            logger.error(f"更新文档失败: {str(e)}")
            return {
                "success": False,
                "error": str(e)
            }
    
    async def delete_document(self, document_id: int) -> Dict[str, Any]:
        """删除文档
        
        功能说明：
        - 通过文档ID获取文档记录；
        - 先从向量存储删除该文档的所有分块向量；
        - 再执行数据库硬删除（调用实例方法 document.delete），依赖外键进行级联删除分块；
        - 最后尝试删除本地存储的原始文件。
        返回结构：{"success": bool, "document_id": int, "vector_result": Any, "timestamp": str}。
        """
        try:
            document = Document.get_by_id(document_id)
            if not document:
                return {
                    "success": False,
                    "error": "文档不存在"
                }
            
            # 从向量存储中删除
            vector_result = await postgresql_vector_service.delete_document(str(document_id))
            
            # 删除文档记录（实例方法，包含标签清理；分块依赖外键级联或上面的向量表清理）
            delete_success = document.delete()
            if not delete_success:
                logger.warning(f"删除文档记录失败: {document_id}")
                return {
                    "success": False,
                    "error": "删除文档记录失败"
                }
            
            # 删除本地文件
            try:
                file_path = os.path.join(self.upload_dir, f"{document.content_hash}.*")
                # 删除匹配的文件
                import glob
                for file in glob.glob(file_path):
                    os.remove(file)
            except Exception as e:
                logger.warning(f"删除本地文件失败: {str(e)}")
            
            result = {
                "success": True,
                "document_id": document_id,
                "vector_result": vector_result,
                "timestamp": datetime.now().isoformat()
            }
            
            logger.info(f"文档 {document_id} 删除成功")
            return result
            
        except Exception as e:
            logger.error(f"删除文档失败: {str(e)}")
            return {
                "success": False,
                "error": str(e)
            }
    
    # ==================== 分类管理 ====================
    
    async def create_category(
        self,
        name: str,
        description: Optional[str] = None,
        parent_id: Optional[int] = None,
        created_by: Optional[int] = None
    ) -> Dict[str, Any]:
        """创建文档分类"""
        try:
            # 检查分类名称是否已存在
            if self._category_name_exists(name, parent_id):
                return {
                    "success": False,
                    "error": "分类名称已存在"
                }
            
            # 如果有父分类，检查父分类是否存在
            if parent_id and not self._category_exists(parent_id):
                return {
                    "success": False,
                    "error": "父分类不存在"
                }
            
            category_id = self._create_category_record(
                name=name,
                description=description,
                parent_id=parent_id,
                created_by=created_by
            )
            
            return {
                "success": True,
                "category_id": category_id,
                "name": name,
                "description": description,
                "parent_id": parent_id,
                "timestamp": datetime.now().isoformat()
            }
            
        except Exception as e:
            logger.error(f"创建分类失败: {str(e)}")
            return {
                "success": False,
                "error": str(e)
            }
    
    async def get_categories(self, include_count: bool = False) -> Dict[str, Any]:
        """获取分类列表"""
        try:
            categories = self._get_all_categories(include_count)
            
            return {
                "categories": categories,
                "total": len(categories),
                "timestamp": datetime.now().isoformat()
            }
            
        except Exception as e:
            logger.error(f"获取分类列表失败: {str(e)}")
            return {
                "categories": [],
                "total": 0,
                "error": str(e)
            }
    
    async def update_category(
        self,
        category_id: int,
        name: Optional[str] = None,
        description: Optional[str] = None,
        parent_id: Optional[int] = None
    ) -> Dict[str, Any]:
        """更新分类信息"""
        try:
            if not self._category_exists(category_id):
                return {
                    "success": False,
                    "error": "分类不存在"
                }
            
            # 检查名称冲突
            if name and self._category_name_exists(name, parent_id, exclude_id=category_id):
                return {
                    "success": False,
                    "error": "分类名称已存在"
                }
            
            # 检查父分类循环引用
            if parent_id and self._would_create_cycle(category_id, parent_id):
                return {
                    "success": False,
                    "error": "不能将分类设置为自己的子分类"
                }
            
            success = self._update_category_record(
                category_id=category_id,
                name=name,
                description=description,
                parent_id=parent_id
            )
            
            if success:
                return {
                    "success": True,
                    "category_id": category_id,
                    "timestamp": datetime.now().isoformat()
                }
            else:
                return {
                    "success": False,
                    "error": "更新失败"
                }
            
        except Exception as e:
            logger.error(f"更新分类失败: {str(e)}")
            return {
                "success": False,
                "error": str(e)
            }
    
    async def delete_category(
        self,
        category_id: int,
        force: bool = False
    ) -> Dict[str, Any]:
        """删除分类"""
        try:
            if not self._category_exists(category_id):
                return {
                    "success": False,
                    "error": "分类不存在"
                }
            
            # 检查是否有子分类
            subcategories = self._get_subcategories(category_id)
            if subcategories and not force:
                return {
                    "success": False,
                    "error": "分类下有子分类，请先删除子分类或使用强制删除"
                }
            
            # 检查是否有文档
            document_count = self._get_category_document_count(category_id)
            if document_count > 0 and not force:
                return {
                    "success": False,
                    "error": f"分类下有 {document_count} 个文档，请先移动文档或使用强制删除"
                }
            
            # 执行删除
            if force:
                # 强制删除：删除所有子分类和文档
                self._force_delete_category(category_id)
            else:
                # 普通删除
                self._delete_category_record(category_id)
            
            return {
                "success": True,
                "category_id": category_id,
                "force": force,
                "timestamp": datetime.now().isoformat()
            }
            
        except Exception as e:
            logger.error(f"删除分类失败: {str(e)}")
            return {
                "success": False,
                "error": str(e)
            }
    
    # ==================== 标签管理 ====================
    
    async def create_tag(
        self,
        name: str,
        color: str = "#1890ff",
        description: Optional[str] = None,
        created_by: Optional[int] = None
    ) -> Dict[str, Any]:
        """创建标签"""
        try:
            # 检查标签名称是否已存在
            if self._tag_name_exists(name):
                return {
                    "success": False,
                    "error": "标签名称已存在"
                }
            
            tag_id = self._create_tag_record(
                name=name,
                color=color,
                description=description,
                created_by=created_by
            )
            
            return {
                "success": True,
                "tag_id": tag_id,
                "name": name,
                "color": color,
                "description": description,
                "timestamp": datetime.now().isoformat()
            }
            
        except Exception as e:
            logger.error(f"创建标签失败: {str(e)}")
            return {
                "success": False,
                "error": str(e)
            }
    
    async def get_tags(self, include_count: bool = False) -> Dict[str, Any]:
        """获取标签列表"""
        try:
            tags = self._get_all_tags(include_count)
            
            return {
                "tags": tags,
                "total": len(tags),
                "timestamp": datetime.now().isoformat()
            }
            
        except Exception as e:
            logger.error(f"获取标签列表失败: {str(e)}")
            return {
                "tags": [],
                "total": 0,
                "error": str(e)
            }
    
    async def delete_tag(self, tag_id: int) -> Dict[str, Any]:
        """删除标签"""
        try:
            if not self._tag_exists(tag_id):
                return {
                    "success": False,
                    "error": "标签不存在"
                }
            
            # 删除标签（会自动从文档中移除关联）
            success = self._delete_tag_record(tag_id)
            
            if success:
                return {
                    "success": True,
                    "tag_id": tag_id,
                    "timestamp": datetime.now().isoformat()
                }
            else:
                return {
                    "success": False,
                    "error": "删除失败"
                }
            
        except Exception as e:
            logger.error(f"删除标签失败: {str(e)}")
            return {
                "success": False,
                "error": str(e)
            }
    
    # ==================== 批量操作 ====================
    
    async def batch_operation(
        self,
        document_ids: List[int],
        operation: str,
        category_id: Optional[int] = None,
        tag_ids: Optional[List[int]] = None,
        user_id: Optional[int] = None
    ) -> Dict[str, Any]:
        """批量操作文档"""
        try:
            results = []
            success_count = 0
            
            for doc_id in document_ids:
                try:
                    if operation == "delete":
                        result = await self.delete_document(doc_id)
                    elif operation == "categorize":
                        result = await self.update_document(
                            document_id=doc_id,
                            category_id=category_id
                        )
                    elif operation == "tag":
                        result = await self.update_document(
                            document_id=doc_id,
                            tag_ids=tag_ids
                        )
                    else:
                        result = {
                            "success": False,
                            "error": f"不支持的操作: {operation}"
                        }
                    
                    results.append({
                        "document_id": doc_id,
                        "success": result["success"],
                        "error": result.get("error") if not result["success"] else None
                    })
                    
                    if result["success"]:
                        success_count += 1
                        
                except Exception as e:
                    results.append({
                        "document_id": doc_id,
                        "success": False,
                        "error": str(e)
                    })
            
            return {
                "operation": operation,
                "total_count": len(document_ids),
                "success_count": success_count,
                "failed_count": len(document_ids) - success_count,
                "results": results,
                "timestamp": datetime.now().isoformat()
            }
            
        except Exception as e:
            logger.error(f"批量操作失败: {str(e)}")
            return {
                "operation": operation,
                "total_count": len(document_ids),
                "success_count": 0,
                "failed_count": len(document_ids),
                "error": str(e)
            }
    
    # ==================== 搜索功能 ====================
    
    async def search_documents(
        self, 
        query: str, 
        top_k: int = 5,
        category: Optional[str] = None,
        tags: Optional[List[str]] = None
    ) -> Dict[str, Any]:
        """搜索文档（支持分类和标签筛选）"""
        try:
            # 使用向量存储搜索
            similar_docs = await postgresql_vector_service.search_similar(
                query=query,
                top_k=top_k,
                filters={
                    "category": category,
                    "tags": tags
                }
            )
            
            return {
                "query": query,
                "results": similar_docs,
                "count": len(similar_docs),
                "filters": {
                    "category": category,
                    "tags": tags
                },
                "timestamp": datetime.now().isoformat()
            }
            
        except Exception as e:
            logger.error(f"文档搜索失败: {str(e)}")
            return {
                "query": query,
                "results": [],
                "count": 0,
                "error": str(e)
            }
    
    async def advanced_search(
        self,
        query: Optional[str] = None,
        category_id: Optional[int] = None,
        tag_ids: Optional[List[int]] = None,
        file_type: Optional[str] = None,
        date_from: Optional[str] = None,
        date_to: Optional[str] = None,
        size_min: Optional[int] = None,
        size_max: Optional[int] = None,
        offset: int = 0,
        limit: int = 20
    ) -> Dict[str, Any]:
        """高级搜索文档"""
        try:
            # 构建搜索条件
            search_filters = {}
            if category_id:
                search_filters["category_id"] = category_id
            if tag_ids:
                search_filters["tag_ids"] = tag_ids
            if file_type:
                search_filters["file_type"] = file_type
            if date_from:
                search_filters["date_from"] = date_from
            if date_to:
                search_filters["date_to"] = date_to
            if size_min:
                search_filters["size_min"] = size_min
            if size_max:
                search_filters["size_max"] = size_max
            
            # 执行搜索
            documents, total = Document.advanced_search(
                query=query,
                filters=search_filters,
                offset=offset,
                limit=limit
            )
            
            document_list = []
            for doc in documents:
                doc_dict = doc.to_dict()
                # 添加额外信息
                doc_dict["chunks_count"] = len(DocumentChunk.get_by_document_id(doc.id))
                if hasattr(doc, 'category_id') and doc.category_id:
                    doc_dict["category"] = self._get_category_info(doc.category_id)
                if hasattr(doc, 'tag_ids') and doc.tag_ids:
                    doc_dict["tags"] = self._get_tags_info(doc.tag_ids)
                document_list.append(doc_dict)
            
            return {
                "query": query,
                "documents": document_list,
                "total": total,
                "offset": offset,
                "limit": limit,
                "filters": search_filters,
                "timestamp": datetime.now().isoformat()
            }
            
        except Exception as e:
            logger.error(f"高级搜索失败: {str(e)}")
            return {
                "query": query,
                "documents": [],
                "total": 0,
                "error": str(e)
            }
    
    # ==================== 统计和健康检查 ====================
    
    async def get_stats(self) -> Dict[str, Any]:
        """获取文档统计信息"""
        try:
            # 获取文档统计
            doc_stats = Document.get_stats()
            
            # 获取分类统计
            category_stats = self._get_category_stats()
            
            # 获取标签统计
            tag_stats = self._get_tag_stats()
            
            # 获取向量存储统计
            vector_stats = await postgresql_vector_service.get_stats()
            
            return {
                "documents": doc_stats,
                "categories": category_stats,
                "tags": tag_stats,
                "vector_store": vector_stats,
                "timestamp": datetime.now().isoformat()
            }
            
        except Exception as e:
            logger.error(f"获取统计信息失败: {str(e)}")
            return {"error": str(e)}
    
    async def health_check(self) -> Dict[str, Any]:
        """健康检查"""
        try:
            health_status = {
                "document_service": True,
                "vector_store": False,
                "upload_dir": False,
                "database": False,
                "overall": False,
                "timestamp": datetime.now().isoformat()
            }
            
            # 检查上传目录
            try:
                health_status["upload_dir"] = os.path.exists(self.upload_dir) and os.access(self.upload_dir, os.W_OK)
            except Exception as e:
                logger.error(f"上传目录检查失败: {str(e)}")
            
            # 检查数据库
            try:
                # 简单的数据库连接测试（函数级注释）
                # 为避免模型未实现 get_stats 导致健康检查失败，改为直接执行 SELECT 1
                from database.config import get_db_connection
                with get_db_connection() as conn:
                    with conn.cursor() as cursor:
                        cursor.execute("SELECT 1")
                        _ = cursor.fetchone()
                health_status["database"] = True
            except Exception as e:
                logger.error(f"数据库检查失败: {str(e)}")
            
            # 检查向量存储
            # 注意：postgresql_vector_service.health_check 为同步方法
            vector_health = postgresql_vector_service.health_check()
            health_status["vector_store"] = vector_health.get("overall", False)
            
            health_status["overall"] = all([
                health_status["document_service"],
                health_status["vector_store"],
                health_status["upload_dir"],
                health_status["database"]
            ])
            
            return health_status
            
        except Exception as e:
            logger.error(f"健康检查失败: {str(e)}")
            return {"overall": False, "error": str(e)}
    
    # ==================== 私有辅助方法 ====================
    
    def _validate_file(self, file_content: bytes, filename: str) -> Dict[str, Any]:
        """验证文件
        
        函数说明（中文）：
        - 检查文件大小是否超过上限。
        - 检查文件类型是否在知识库允许的类型范围内。
        - 严格校验扩展名与策略映射（例如 excel 仅支持 .xlsx）。
        - 校验文件非空。
        
        返回：{"valid": bool, "error": str?}
        """
        try:
            # 检查文件大小
            if len(file_content) > self.max_file_size:
                return {
                    "valid": False,
                    "error": f"文件大小超过限制 ({self.max_file_size / 1024 / 1024:.1f}MB)"
                }
            
            # 检查文件类型
            file_type = get_file_type(filename)
            # 先用内部支持类型集合过滤（排除图片/音视频等）
            # 仅允许可向量化类型（排除图片、音视频等）
            if file_type not in self.supported_types:
                return {
                    "valid": False,
                    "error": f"该文件类型目前不支持向量化: {file_type}"
                }
            
            # 严格校验扩展名映射（与前端展示保持一致）
            policy = self.get_upload_policy()
            allowed_exts = policy.get("supported_types", {}).get(file_type, [])
            ext = os.path.splitext(filename)[1].lower()
            if allowed_exts and ext not in allowed_exts:
                # 定制 excel 的报错文案，其他类型统一提示
                if file_type == "excel":
                    return {
                        "valid": False,
                        "error": "Excel 仅支持 .xlsx，请转换后再上传"
                    }
                return {
                    "valid": False,
                    "error": f"{file_type} 仅支持: {', '.join(allowed_exts)}"
                }
            
            # 检查文件内容
            if len(file_content) == 0:
                return {
                    "valid": False,
                    "error": "文件内容为空"
                }
            
            return {"valid": True}
            
        except Exception as e:
            return {
                "valid": False,
                "error": f"文件验证失败: {str(e)}"
            }

    async def _cleanup_failed_upload(self, document: Document):
        """失败上传清理函数
        
        函数说明（中文）：
        - 删除向量存储中的该文档条目。
        - 删除数据库中的文档记录（含标签关联，分块依赖外键级联删除）。
        - 删除本地存储的原始文件（基于 content_hash 通配符）。
        """
        try:
            # 删除向量数据
            try:
                await postgresql_vector_service.delete_document(str(document.id))
            except Exception:
                pass
            # 删除数据库记录
            try:
                document.delete()
            except Exception:
                pass
            # 删除本地文件
            try:
                import glob
                for f in glob.glob(os.path.join(self.upload_dir, f"{document.content_hash}.*")):
                    try:
                        os.remove(f)
                    except Exception:
                        pass
            except Exception:
                pass
        except Exception as e:
            logger.debug(f"失败清理出现异常: {e}")
    
    async def _save_file(self, file_content: bytes, filename: str, content_hash: str) -> str:
        """保存文件到本地"""
        try:
            # 生成文件路径
            file_ext = os.path.splitext(filename)[1]
            safe_filename = f"{content_hash}{file_ext}"
            file_path = os.path.join(self.upload_dir, safe_filename)
            
            # 保存文件
            with open(file_path, 'wb') as f:
                f.write(file_content)
            
            logger.info(f"文件保存到: {file_path}")
            return file_path
            
        except Exception as e:
            logger.error(f"文件保存失败: {str(e)}")
            raise
    
    async def _extract_text(self, file_content: bytes, file_type: str) -> str:
        """提取文档文本内容"""
        try:
            if file_type == 'pdf':
                return await self._extract_pdf_text(file_content)
            elif file_type == 'word':
                return await self._extract_word_text(file_content)
            elif file_type in ['text', 'markdown']:
                return file_content.decode('utf-8', errors='ignore')
            elif file_type == 'html':
                return await self._extract_html_text(file_content)
            elif file_type == 'powerpoint':
                return await self._extract_powerpoint_text(file_content)
            elif file_type == 'excel':
                return await self._extract_excel_text(file_content)
            else:
                raise ValueError(f"不支持的文件类型: {file_type}")
                
        except Exception as e:
            logger.error(f"文本提取失败: {str(e)}")
            return ""
    
    async def _extract_pdf_text(self, file_content: bytes) -> str:
        """提取PDF文本"""
        try:
            pdf_file = io.BytesIO(file_content)
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            
            text_content = ""
            for page in pdf_reader.pages:
                text_content += page.extract_text() + "\n"
            
            return text_content.strip()
            
        except Exception as e:
            logger.error(f"PDF文本提取失败: {str(e)}")
            return ""
    
    async def _extract_word_text(self, file_content: bytes) -> str:
        """提取Word文档文本"""
        try:
            doc_file = io.BytesIO(file_content)
            doc = DocxDocument(doc_file)
            
            text_content = ""
            for paragraph in doc.paragraphs:
                text_content += paragraph.text + "\n"
            
            return text_content.strip()
            
        except Exception as e:
            logger.error(f"Word文档文本提取失败: {str(e)}")
            return ""
    
    async def _extract_html_text(self, file_content: bytes) -> str:
        """提取HTML文本"""
        try:
            # 简单的HTML标签移除
            html_content = file_content.decode('utf-8', errors='ignore')
            
            # 移除HTML标签（简单实现）
            import re
            text_content = re.sub(r'<[^>]+>', '', html_content)
            text_content = re.sub(r'\s+', ' ', text_content)
            
            return text_content.strip()
            
        except Exception as e:
            logger.error(f"HTML文本提取失败: {str(e)}")
            return ""

    async def _extract_powerpoint_text(self, file_content: bytes) -> str:
        """提取PPTX演示文稿文本
        
        - 仅支持 Office Open XML 格式（.pptx）。
        - 若为旧版 .ppt 格式或解析库缺失，返回空字符串并记录日志。
        """
        try:
            # 延迟导入，避免环境未安装导致模块导入阶段失败
            try:
                from pptx import Presentation
            except ImportError:
                Presentation = None  # 使用ZIP解析回退

            texts: List[str] = []

            if Presentation is not None:
                ppt_file = io.BytesIO(file_content)
                prs = Presentation(ppt_file)

                for slide_idx, slide in enumerate(prs.slides):
                    # 提取文本框内容
                    for shape in slide.shapes:
                        try:
                            # 文本框/占位符文本
                            if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                                # shape.text 会聚合段落文本
                                if getattr(shape, 'text', None):
                                    texts.append(shape.text)
                                    continue
                                # 兜底遍历段落与runs
                                try:
                                    tf = shape.text_frame
                                    for p in getattr(tf, 'paragraphs', []):
                                        # 聚合段落全文
                                        paragraph_text = ''.join([run.text for run in getattr(p, 'runs', [])])
                                        if paragraph_text:
                                            texts.append(paragraph_text)
                                except Exception:
                                    pass

                            # 表格中的文本
                            if hasattr(shape, 'has_table') and shape.has_table:
                                table = shape.table
                                for row in table.rows:
                                    for cell in row.cells:
                                        try:
                                            cell_text = getattr(cell, 'text', '')
                                            if cell_text:
                                                texts.append(cell_text)
                                        except Exception:
                                            continue
                        except Exception as inner_e:
                            logger.debug(f"PPTX第{slide_idx+1}页形状解析警告: {inner_e}")
            else:
                # 无 python-pptx 时，回退使用 zipfile + XML 解析 a:t 文本节点
                import zipfile
                from xml.etree import ElementTree as ET

                try:
                    zf = zipfile.ZipFile(io.BytesIO(file_content))
                except Exception as ze:
                    logger.error(f"PPTX压缩结构解析失败: {ze}")
                    return ""

                slide_xml_files = [name for name in zf.namelist() if name.startswith('ppt/slides/slide') and name.endswith('.xml')]
                for name in slide_xml_files:
                    try:
                        with zf.open(name) as f:
                            xml_bytes = f.read()
                        # 去除命名空间以便简单匹配
                        try:
                            root = ET.fromstring(xml_bytes)
                        except Exception as xe:
                            logger.debug(f"解析XML失败 {name}: {xe}")
                            continue

                        # 遍历所有文本节点（命名空间通常为 a:t）
                        for node in root.iter():
                            tag = node.tag
                            if isinstance(tag, str) and tag.endswith('}t'):
                                val = (node.text or '').strip()
                                if val:
                                    texts.append(val)
                    except Exception as se:
                        logger.debug(f"读取PPTX幻灯片失败 {name}: {se}")

            content = '\n'.join([t.strip() for t in texts if t and t.strip()])
            return content.strip()
        except Exception as e:
            # 对于老版 .ppt 或损坏文件解析失败
            logger.error(f"PPTX文本提取失败: {str(e)}")
            return ""

    async def _extract_excel_text(self, file_content: bytes) -> str:
        """提取Excel工作簿文本
        
        - 依赖 openpyxl 读取 .xlsx 内容（只读遍历所有工作表与单元格）。
        - 非空行以制表符连接，工作表之间增加标题分隔。
        """
        try:
            try:
                from openpyxl import load_workbook
            except ImportError:
                logger.error("Excel文本提取失败: openpyxl 未安装")
                return ""

            wb_file = io.BytesIO(file_content)
            workbook = load_workbook(wb_file, read_only=True, data_only=True)
            parts: List[str] = []

            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                row_texts: List[str] = []
                for row in sheet.iter_rows(values_only=True):
                    cells = ["" if (c is None) else str(c) for c in row]
                    # 跳过全空行
                    if any(cell.strip() for cell in cells):
                        row_texts.append("\t".join(cells))
                if row_texts:
                    parts.append(f"=== 工作表: {sheet_name} ===\n" + "\n".join(row_texts))

            content = "\n\n".join(parts).strip()
            return content
        except Exception as e:
            logger.error(f"Excel文本提取失败: {str(e)}")
            return ""
    
    async def _save_chunks(self, document_id: int, chunks: List[str]):
        """保存文档分块到数据库"""
        try:
            # 导入数据库模型
            from database.models import DocumentChunk as DBDocumentChunk
            
            # 创建文档块对象列表
            chunk_objects = []
            for i, chunk in enumerate(chunks):
                chunk_obj = DBDocumentChunk(
                    document_id=document_id,
                    chunk_index=i,
                    content=chunk,
                    content_length=len(chunk),
                    metadata={"length": len(chunk)}
                )
                chunk_objects.append(chunk_obj)
                logger.debug(f"准备保存分块 {i}: {len(chunk)} 字符")
            
            # 批量保存到数据库
            if chunk_objects:
                created_chunks = await DBDocumentChunk.create_batch(chunk_objects)
                logger.info(f"文档 {document_id} 的 {len(created_chunks)} 个分块保存完成")
            else:
                logger.warning(f"文档 {document_id} 没有分块需要保存")
            
        except Exception as e:
            logger.error(f"保存文档分块失败: {str(e)}")
            raise

    # ==================== 支持类型/策略查询 ====================
    def get_upload_policy(self) -> Dict[str, Any]:
        """返回知识库上传策略与支持的文件类型
        
        - 包含最大文件大小（MB）
        - 返回可向量化类型及其扩展名列表
        """
        return {
            "max_file_size_mb": int(self.max_file_size / 1024 / 1024),
            "supported_types": {
                "pdf": [".pdf"],
                "word": [".doc", ".docx"],
                "excel": [".xlsx"],
                "powerpoint": [".ppt", ".pptx"],
                "text": [".txt"],
                "markdown": [".md"],
                "html": [".htm", ".html"],
            }
        }
    
    # 分类相关辅助方法
    def _category_name_exists(self, name: str, parent_id: Optional[int] = None, exclude_id: Optional[int] = None) -> bool:
        """检查分类名称是否存在"""
        # 这里应该查询数据库
        # 暂时返回False，实际实现需要查询categories表
        return False
    
    def _category_exists(self, category_id: int) -> bool:
        """检查分类是否存在"""
        # 这里应该查询数据库
        return True
    
    def _create_category_record(self, name: str, description: Optional[str], parent_id: Optional[int], created_by: Optional[int]) -> int:
        """创建分类记录"""
        # 这里应该插入数据库并返回ID
        return 1
    
    def _get_all_categories(self, include_count: bool = False) -> List[Dict[str, Any]]:
        """获取所有分类"""
        # 这里应该查询数据库
        return []
    
    def _update_category_record(self, category_id: int, name: Optional[str], description: Optional[str], parent_id: Optional[int]) -> bool:
        """更新分类记录"""
        # 这里应该更新数据库
        return True
    
    def _would_create_cycle(self, category_id: int, parent_id: int) -> bool:
        """检查是否会创建循环引用"""
        # 这里应该检查分类层级关系
        return False
    
    def _get_subcategories(self, category_id: int) -> List[int]:
        """获取子分类"""
        # 这里应该查询数据库
        return []
    
    def _get_category_document_count(self, category_id: int) -> int:
        """获取分类下的文档数量"""
        # 这里应该查询数据库
        return 0
    
    def _force_delete_category(self, category_id: int):
        """强制删除分类（包括子分类和文档）"""
        # 这里应该递归删除所有子分类和文档
        pass
    
    def _delete_category_record(self, category_id: int):
        """删除分类记录"""
        # 这里应该删除数据库记录
        pass
    
    def _get_category_info(self, category_id: int) -> Optional[Dict[str, Any]]:
        """获取分类信息"""
        # 这里应该查询数据库
        return None
    
    def _get_category_stats(self) -> Dict[str, Any]:
        """获取分类统计"""
        return {
            "total_categories": 0,
            "categories_with_documents": 0
        }
    
    # 标签相关辅助方法
    def _tag_name_exists(self, name: str) -> bool:
        """检查标签名称是否存在"""
        # 这里应该查询数据库
        return False
    
    def _tag_exists(self, tag_id: int) -> bool:
        """检查标签是否存在"""
        # 这里应该查询数据库
        return True
    
    def _create_tag_record(self, name: str, color: str, description: Optional[str], created_by: Optional[int]) -> int:
        """创建标签记录"""
        # 这里应该插入数据库并返回ID
        return 1
    
    def _get_all_tags(self, include_count: bool = False) -> List[Dict[str, Any]]:
        """获取所有标签"""
        # 这里应该查询数据库
        return []
    
    def _delete_tag_record(self, tag_id: int) -> bool:
        """删除标签记录"""
        # 这里应该删除数据库记录
        return True
    
    def _get_tags_info(self, tag_ids: List[int]) -> List[Dict[str, Any]]:
        """获取标签信息"""
        # 这里应该查询数据库
        return []
    
    def _get_tag_stats(self) -> Dict[str, Any]:
        """获取标签统计"""
        return {
            "total_tags": 0,
            "tags_in_use": 0
        }
    
    # 分类管理方法
    def create_category(self, name: str, description: str = None, color: str = None) -> Dict[str, Any]:
        """创建分类"""
        try:
            from app.models.document import Category
            
            # 检查分类名是否已存在
            if self._category_name_exists(name):
                raise HTTPException(status_code=400, detail="分类名已存在")
            
            category = Category.create_category(name=name, description=description, color=color)
            if not category:
                raise HTTPException(status_code=500, detail="创建分类失败")
            
            return category.to_dict()
            
        except HTTPException:
            raise
        except Exception as e:
            logger.error(f"创建分类失败: {str(e)}")
            raise HTTPException(status_code=500, detail="创建分类失败")
    
    def get_categories(self) -> List[Dict[str, Any]]:
        """获取所有分类"""
        try:
            from app.models.document import Category
            categories = Category.list_all()
            return [category.to_dict() for category in categories]
        except Exception as e:
            logger.error(f"获取分类列表失败: {str(e)}")
            raise HTTPException(status_code=500, detail="获取分类列表失败")
    
    def update_category(self, category_id: int, name: str = None, 
                       description: str = None, color: str = None) -> Dict[str, Any]:
        """更新分类"""
        try:
            from app.models.document import Category
            
            # 检查分类是否存在
            category = Category.get_by_id(category_id)
            if not category:
                raise HTTPException(status_code=404, detail="分类不存在")
            
            # 如果更新名称，检查是否与其他分类重名
            if name and name != category.name and self._category_name_exists(name):
                raise HTTPException(status_code=400, detail="分类名已存在")
            
            success = Category.update_category(
                category_id=category_id, name=name, 
                description=description, color=color
            )
            
            if not success:
                raise HTTPException(status_code=500, detail="更新分类失败")
            
            # 返回更新后的分类信息
            updated_category = Category.get_by_id(category_id)
            return updated_category.to_dict()
            
        except HTTPException:
            raise
        except Exception as e:
            logger.error(f"更新分类失败: {str(e)}")
            raise HTTPException(status_code=500, detail="更新分类失败")
    
    def delete_category(self, category_id: int) -> Dict[str, str]:
        """删除分类"""
        try:
            from app.models.document import Category
            
            # 检查分类是否存在
            category = Category.get_by_id(category_id)
            if not category:
                raise HTTPException(status_code=404, detail="分类不存在")
            
            # 检查是否有文档使用此分类
            documents, _ = Document.list_documents(category_id=category_id, limit=1)
            if documents:
                raise HTTPException(status_code=400, detail="该分类下还有文档，无法删除")
            
            success = Category.delete_category(category_id)
            if not success:
                raise HTTPException(status_code=500, detail="删除分类失败")
            
            return {"message": "分类删除成功"}
            
        except HTTPException:
            raise
        except Exception as e:
            logger.error(f"删除分类失败: {str(e)}")
            raise HTTPException(status_code=500, detail="删除分类失败")
    
    # 标签管理方法
    def create_tag(self, name: str, color: str = None) -> Dict[str, Any]:
        """创建标签"""
        try:
            from app.models.document import Tag
            
            # 检查标签名是否已存在
            if self._tag_name_exists(name):
                raise HTTPException(status_code=400, detail="标签名已存在")
            
            tag = Tag.create_tag(name=name, color=color)
            if not tag:
                raise HTTPException(status_code=500, detail="创建标签失败")
            
            return tag.to_dict()
            
        except HTTPException:
            raise
        except Exception as e:
            logger.error(f"创建标签失败: {str(e)}")
            raise HTTPException(status_code=500, detail="创建标签失败")
    
    def get_tags(self) -> List[Dict[str, Any]]:
        """获取所有标签"""
        try:
            from app.models.document import Tag
            tags = Tag.list_all()
            return [tag.to_dict() for tag in tags]
        except Exception as e:
            logger.error(f"获取标签列表失败: {str(e)}")
            raise HTTPException(status_code=500, detail="获取标签列表失败")
    
    def delete_tag(self, tag_id: int) -> Dict[str, str]:
        """删除标签"""
        try:
            from app.models.document import Tag
            
            # 检查标签是否存在
            tag = Tag.get_by_id(tag_id)
            if not tag:
                raise HTTPException(status_code=404, detail="标签不存在")
            
            success = Tag.delete_tag(tag_id)
            if not success:
                raise HTTPException(status_code=500, detail="删除标签失败")
            
            return {"message": "标签删除成功"}
            
        except HTTPException:
            raise
        except Exception as e:
            logger.error(f"删除标签失败: {str(e)}")
            raise HTTPException(status_code=500, detail="删除标签失败")
    
    # 私有辅助方法
    def _category_name_exists(self, name: str) -> bool:
        """检查分类名是否已存在"""
        try:
            from app.models.document import Category
            categories = Category.list_all()
            return any(cat.name == name for cat in categories)
        except Exception:
            return False
    
    def _tag_name_exists(self, name: str) -> bool:
        """检查标签名是否已存在"""
        try:
            from app.models.document import Tag
            tags = Tag.list_all()
            return any(tag.name == name for tag in tags)
        except Exception:
            return False
    
    # 高级搜索方法
    def advanced_search(self, 
                       search_term: str = None,
                       category_id: int = None,
                       tag_ids: List[int] = None,
                       file_type: str = None,
                       status: str = None,
                       date_from: str = None,
                       date_to: str = None,
                       page: int = 1,
                       limit: int = 20) -> Dict[str, Any]:
        """高级搜索文档"""
        try:
            # 构建搜索参数
            search_params = {
                'search_term': search_term,
                'category_id': category_id,
                'tag_ids': tag_ids,
                'file_type': file_type,
                'status': status,
                'page': page,
                'limit': limit
            }
            
            # 获取文档列表
            documents, total_count = Document.list_documents(**search_params)
            
            # 如果有日期范围过滤，进一步筛选
            if date_from or date_to:
                filtered_docs = []
                for doc in documents:
                    doc_date = doc.created_at.date() if hasattr(doc.created_at, 'date') else doc.created_at
                    
                    # 检查日期范围
                    if date_from:
                        from datetime import datetime
                        from_date = datetime.strptime(date_from, '%Y-%m-%d').date()
                        if doc_date < from_date:
                            continue
                    
                    if date_to:
                        from datetime import datetime
                        to_date = datetime.strptime(date_to, '%Y-%m-%d').date()
                        if doc_date > to_date:
                            continue
                    
                    filtered_docs.append(doc)
                
                documents = filtered_docs
                total_count = len(filtered_docs)
            
            # 添加分类信息
            result_docs = []
            for doc in documents:
                doc_dict = doc.to_dict()
                
                # 添加分类信息
                if doc.category_id:
                    from app.models.document import Category
                    category = Category.get_by_id(doc.category_id)
                    doc_dict['category'] = category.to_dict() if category else None
                else:
                    doc_dict['category'] = None
                
                result_docs.append(doc_dict)
            
            # 计算分页信息
            total_pages = (total_count + limit - 1) // limit
            
            return {
                "documents": result_docs,
                "pagination": {
                    "current_page": page,
                    "total_pages": total_pages,
                    "total_count": total_count,
                    "page_size": limit
                }
            }
            
        except Exception as e:
            logger.error(f"高级搜索失败: {str(e)}")
            raise HTTPException(status_code=500, detail="搜索失败")

# 全局文档服务实例
document_service = DocumentService()
